# LocalRAG_Pro\app\build_index.py
# -*- coding: utf-8 -*-
#
# ドキュメントをベクトルインデックスに変換するスクリプト
# 対応形式: PDF, XLSX, DOCX, PPTX, CSV, TXT
# 使い方: py .\app\build_index.py

import csv
import os
import re
import sys
import warnings
from collections import defaultdict
from pathlib import Path
from typing import Dict, List, Set, Tuple

from llama_index.core import Document, Settings, VectorStoreIndex
from llama_index.core.node_parser import SentenceSplitter
from llama_index.embeddings.huggingface import HuggingFaceEmbedding

# -------------------------------------------------------
# Windows cp932 コンソールでの UnicodeEncodeError を防ぐ
# -------------------------------------------------------
sys.stdout.reconfigure(encoding="utf-8", errors="replace")
sys.stderr.reconfigure(encoding="utf-8", errors="replace")

# PIL の超大画像警告を抑制（ポスター等の巨大画像PDF用）
warnings.filterwarnings("ignore", message=".*DecompressionBomb.*")


# =========================
# Config
# =========================

SUPPORTED_EXTS = {".pdf", ".docx", ".xlsx", ".pptx", ".csv", ".txt"}

# インデックス作成時に常に除外するフォルダ名
ALWAYS_EXCLUDE_DIR_NAMES = {
    ".venv", "venv", "__pycache__", ".git", "node_modules",
    ".cursor", ".idea", ".vscode", "templates",
}

# 環境変数キー
PDF_LIMIT_ENV          = "PDF_LIMIT"           # PDF件数上限（テスト用） 例: $env:PDF_LIMIT="50"
ENABLE_OCR_ENV         = "ENABLE_OCR"          # OCR有効/無効  例: $env:ENABLE_OCR="0"
OCR_LANG_ENV           = "OCR_LANG"            # OCR言語        例: $env:OCR_LANG="eng+jpn"
OCR_MIN_TEXT_CHARS_ENV = "OCR_MIN_TEXT_CHARS"  # OCR閾値        例: $env:OCR_MIN_TEXT_CHARS="50"
OCR_CACHE_ENV          = "OCR_CACHE"           # OCRキャッシュ  例: $env:OCR_CACHE="1"
SHOW_ALL_SKIP_ENV      = "SHOW_ALL_SKIP"       # スキップ全件表示 例: $env:SHOW_ALL_SKIP="1"
CHUNK_SIZE_ENV         = "CHUNK_SIZE"          # チャンクサイズ  例: $env:CHUNK_SIZE="512"
CHUNK_OVERLAP_ENV      = "CHUNK_OVERLAP"       # チャンク重複    例: $env:CHUNK_OVERLAP="200"


# =========================
# Paths
# =========================

def get_paths() -> Tuple[Path, Path]:
    script_dir  = Path(__file__).resolve().parent
    default_root = script_dir.parent  # LocalRAG_Pro フォルダ

    root_dir     = Path(os.environ.get("US_ADMIN_ROOT", str(default_root))).resolve()
    default_index = Path.home() / "Documents" / "LocalRAG_Pro_index" / "index"
    index_dir    = Path(os.environ.get("US_ADMIN_INDEX_DIR", str(default_index))).resolve()

    return root_dir, index_dir


def file_metadata_fn(file_path: str) -> dict:
    p = Path(file_path)
    try:
        stat = p.stat()
        mtime, size = stat.st_mtime, stat.st_size
    except Exception:
        mtime, size = None, None
    return {
        "file_path": str(p.resolve()),
        "file_name": p.name,
        "ext":       p.suffix.lower(),
        "mtime":     mtime,
        "size":      size,
    }


# =========================
# OCR helpers
# =========================

def _ocr_enabled() -> bool:
    return os.environ.get(ENABLE_OCR_ENV, "1").strip() not in ("0", "false", "False", "no", "NO")

def _ocr_lang() -> str:
    return os.environ.get(OCR_LANG_ENV, "eng+jpn").strip() or "eng+jpn"

def _ocr_min_text_chars() -> int:
    v = os.environ.get(OCR_MIN_TEXT_CHARS_ENV, "50").strip()
    return int(v) if v.isdigit() else 50

def _ocr_cache_enabled() -> bool:
    return os.environ.get(OCR_CACHE_ENV, "1").strip() not in ("0", "false", "False", "no", "NO")

def _safe_filename(s: str) -> str:
    return re.sub(r"[^A-Za-z0-9._-]+", "_", s)


# =========================
# Readers
# =========================

def read_txt(path: Path) -> str:
    # utf-8-sig を最初に試すことで BOM を自動除去する
    for enc in ("utf-8-sig", "utf-8", "cp932", "utf-16"):
        try:
            return path.read_text(encoding=enc, errors="ignore").strip()
        except Exception:
            continue
    return ""


def read_csv(path: Path) -> str:
    rows: List[str] = []
    for enc in ("utf-8", "utf-8-sig", "cp932", "utf-16"):
        try:
            with path.open(encoding=enc, newline="", errors="ignore") as f:
                reader = csv.reader(f)
                for row in reader:
                    vals = [v.strip() for v in row if v.strip()]
                    if vals:
                        rows.append(" | ".join(vals))
            return "\n".join(rows).strip()
        except Exception:
            continue
    return ""


def read_docx(path: Path) -> str:
    try:
        from docx import Document as DocxDocument
        d = DocxDocument(str(path))
        paras = [p.text.strip() for p in d.paragraphs if p.text.strip()]
        for table in d.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    paras.append(" | ".join(cells))
        return "\n".join(paras).strip()
    except Exception:
        try:
            import docx2txt
            return (docx2txt.process(str(path)) or "").strip()
        except Exception:
            return ""


def read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs    = Presentation(str(path))
        slides: List[str] = []
        for i, slide in enumerate(prs.slides, 1):
            texts: List[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = "".join(run.text for run in para.runs).strip()
                        if t:
                            texts.append(t)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            texts.append(" | ".join(cells))
            if texts:
                slides.append(f"[SLIDE {i}]\n" + "\n".join(texts))
        return "\n\n".join(slides).strip()
    except Exception:
        return ""


def read_xlsx(path: Path) -> str:
    try:
        import openpyxl
        wb       = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
        max_rows = int(os.environ.get("XLSX_MAX_ROWS", "5000"))
        max_cols = int(os.environ.get("XLSX_MAX_COLS", "50"))
        chunks: List[str] = []
        for ws in wb.worksheets:
            rows_text: List[str] = []
            for rcount, row in enumerate(ws.iter_rows(values_only=True), 1):
                if rcount > max_rows:
                    break
                vals = [str(v).strip() for v in row[:max_cols] if v is not None and str(v).strip()]
                if vals:
                    rows_text.append(" | ".join(vals))
            if rows_text:
                chunks.append(f"[SHEET] {ws.title}\n" + "\n".join(rows_text))
        try:
            wb.close()
        except Exception:
            pass
        return "\n\n".join(chunks).strip()
    except Exception:
        return ""


def read_pdf_text(path: Path) -> str:
    try:
        import fitz  # pymupdf
    except ImportError as exc:
        raise ImportError("PyMuPDF が未インストールです。pip install pymupdf を実行してください。") from exc

    doc = fitz.open(str(path))
    try:
        texts: List[str] = []
        for i in range(len(doc)):
            t = doc.load_page(i).get_text("text") or ""
            if t.strip():
                texts.append(t)
        return "\n".join(texts).strip()
    finally:
        doc.close()


def ocr_pdf_to_text(pdf_path: Path, cache_dir: Path) -> str:
    from pdf2image import convert_from_path
    import pytesseract

    lang       = _ocr_lang()
    cache_key  = f"{pdf_path.name}_{pdf_path.stat().st_size}_{int(pdf_path.stat().st_mtime)}"
    cache_file = cache_dir / (_safe_filename(cache_key) + ".ocr.txt")

    if _ocr_cache_enabled() and cache_file.exists():
        return read_txt(cache_file)

    images = convert_from_path(str(pdf_path), dpi=200)
    texts  = [t for img in images if (t := pytesseract.image_to_string(img, lang=lang).strip())]
    out    = "\n\n".join(texts).strip()

    if _ocr_cache_enabled():
        cache_dir.mkdir(parents=True, exist_ok=True)
        try:
            cache_file.write_text(out, encoding="utf-8")
        except Exception:
            pass

    return out


# =========================
# Collection
# =========================

def collect_documents(root_dir: Path, exclude_dirs: Set[str], index_dir: Path) -> List[Document]:
    docs: List[Document]   = []
    skipped: List[str]     = []
    counts: Dict[str, int] = defaultdict(int)

    pdf_limit = None
    limit_str = os.environ.get(PDF_LIMIT_ENV, "").strip()
    if limit_str.isdigit():
        pdf_limit = int(limit_str)

    ocr_min_chars = _ocr_min_text_chars()
    do_ocr        = _ocr_enabled()
    ocr_cache_dir = index_dir / "_ocr_cache"
    pdf_seen      = 0

    # 対象ファイルを事前収集してから進捗表示
    candidates: List[Path] = []
    for p in root_dir.rglob("*"):
        if not p.is_file():
            continue
        if p.suffix.lower() not in SUPPORTED_EXTS:
            continue
        sp = str(p.resolve())
        if any(sp.startswith(ex) for ex in exclude_dirs):
            continue
        if any(part in ALWAYS_EXCLUDE_DIR_NAMES for part in p.parts):
            continue
        candidates.append(p)

    total = len(candidates)
    print(f"[INFO] 対象ファイル数: {total} 件")

    for idx, p in enumerate(candidates, 1):
        ext = p.suffix.lower()
        sp  = str(p.resolve())

        # 進捗表示（日本語ファイル名を安全に表示）
        try:
            print(f"  [{idx:>4}/{total}] {p.name}", end="\r", flush=True)
        except Exception:
            print(f"  [{idx:>4}/{total}] (表示不可)", end="\r", flush=True)

        try:
            text = ""

            if ext == ".txt":
                text = read_txt(p)

            elif ext == ".csv":
                text = read_csv(p)

            elif ext == ".docx":
                text = read_docx(p)

            elif ext == ".pptx":
                text = read_pptx(p)

            elif ext == ".xlsx":
                text = read_xlsx(p)

            elif ext == ".pdf":
                pdf_seen += 1
                if pdf_limit is not None and pdf_seen > pdf_limit:
                    continue
                text = read_pdf_text(p)
                if do_ocr and len(text.strip()) < ocr_min_chars:
                    try:
                        ocr_text = ocr_pdf_to_text(p, cache_dir=ocr_cache_dir)
                        if ocr_text and len(ocr_text.strip()) >= ocr_min_chars:
                            text = ocr_text
                    except Exception as e:
                        skipped.append(f"[OCR-SKIP] {sp} :: {e}")

            text = (text or "").strip()
            if not text:
                skipped.append(f"[EMPTY] {sp}")
                continue

            docs.append(Document(text=text, metadata=file_metadata_fn(sp)))
            counts[ext] += 1

        except Exception as e:
            skipped.append(f"[ERROR] {sp} :: {e}")

    print()  # 進捗行の改行

    # スキップ表示
    show_all = os.environ.get(SHOW_ALL_SKIP_ENV, "0").strip() not in ("0", "false", "no")
    limit    = len(skipped) if show_all else 15
    if skipped:
        print(f"[INFO] スキップ: {len(skipped)} 件 (表示: {min(limit, len(skipped))} 件 / 全件表示: $env:SHOW_ALL_SKIP=\"1\")")
        for s in skipped[:limit]:
            print("  ", s)

    # ファイル種別ごとの集計
    print("\n[INFO] 読み込み成功 内訳:")
    for ext in sorted(counts):
        print(f"       {ext:<8} {counts[ext]:>4} 件")

    return docs


# =========================
# Main
# =========================

def main():
    root_dir, index_dir = get_paths()

    if not root_dir.exists():
        raise FileNotFoundError(f"[ERROR] フォルダが見つかりません: {root_dir}")
    index_dir.mkdir(parents=True, exist_ok=True)

    chunk_size    = int(os.environ.get(CHUNK_SIZE_ENV, "512"))
    chunk_overlap = int(os.environ.get(CHUNK_OVERLAP_ENV, "200"))

    print("=" * 55)
    print("  LocalRAG_Pro インデックス作成")
    print("=" * 55)
    print(f"[INFO] 読込元  : {root_dir}")
    print(f"[INFO] 保存先  : {index_dir}")
    print(f"[INFO] OCR     : {'有効' if _ocr_enabled() else '無効'}  言語={_ocr_lang()}  閾値={_ocr_min_text_chars()}文字")
    print(f"[INFO] チャンク: size={chunk_size}  overlap={chunk_overlap}")
    print("-" * 55)

    # Embedding モデル（build と ui で必ず同じモデルを使う）
    # multilingual-e5-large: small比で精度大幅向上（初回DL約500MB）
    Settings.embed_model = HuggingFaceEmbedding(model_name="intfloat/multilingual-e5-large")
    Settings.node_parser = SentenceSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)

    exclude_dirs = {str(index_dir)}

    docs = collect_documents(root_dir=root_dir, exclude_dirs=exclude_dirs, index_dir=index_dir)
    print(f"\n[INFO] ドキュメント合計: {len(docs)} 件")

    if not docs:
        print("[WARN] 読み込めたドキュメントが0件です。フォルダを確認してください。")
        return

    print("\n[INFO] ベクトルインデックスを構築中... (数分かかります)")
    index = VectorStoreIndex.from_documents(docs, show_progress=True)
    index.storage_context.persist(persist_dir=str(index_dir))

    print("\n" + "=" * 55)
    print(f"  ✅ インデックス作成完了！  {len(docs)} 件")
    print(f"     保存先: {index_dir}")
    print("=" * 55)
    print("  次のコマンドでAPIサーバーを起動:")
    print("  python app/api_server.py")
    print("=" * 55)


if __name__ == "__main__":
    main()
